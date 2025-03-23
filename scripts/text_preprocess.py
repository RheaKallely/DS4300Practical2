import os
import re
import fitz  
from pptx import Presentation
import nltk
from transformers import AutoTokenizer
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
nltk.download('punkt')
nltk.download('stopwords')

def extract_text_from_pdf(pdf_path):
    """Extracts text from a PDF file."""
    text = ""
    try:
        with fitz.open(pdf_path) as doc:
            for page in doc:
                text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"Error extracting {pdf_path}: {e}")
    return text

def extract_text_from_pptx(pptx_path):
    """Extracts text from a PPTX file."""
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting {pptx_path}: {e}")
    return text

def preprocess_text(text):
    """
    Cleans and preprocesses text by:
    - Converting to lowercase
    - Removing extra whitespace
    - Removing punctuation
    - Tokenizing
    - Removing stopwords
    """
    text = text.lower()
    text = re.sub(r'\s+', ' ', text).strip()  
    text = re.sub(r'[^\w\s]', '', text)  
    
    tokens = word_tokenize(text)
    stop_words = set(stopwords.words('english'))
    tokens = [word for word in tokens if word not in stop_words]

    return ' '.join(tokens)

tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
def chunk_text(text, chunk_size=500, overlap=50):
    """
    Splits text into chunks using a model-specific tokenizer.
    """
    tokens = tokenizer.encode(text, add_special_tokens=False)
    chunks = []
    
    for i in range(0, len(tokens), chunk_size - overlap):
        chunk = tokens[i:i + chunk_size]
        chunks.append(tokenizer.decode(chunk))
    
    return chunks

def process_notes_folder(folder_path="data/notes"):
    """
    Processes all PDFs and PPTX files in the 'notes' folder.
    Extracts, preprocesses, and chunks text from each file.
    """
    if not os.path.exists(folder_path):
        print(f"Error: The folder '{folder_path}' does not exist!")
        return []
    
    all_chunks = []
    
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        
        if filename.endswith(".pdf"):
            print(f"Processing PDF: {filename}")
            text = extract_text_from_pdf(file_path)
        
        elif filename.endswith(".pptx"):
            print(f"Processing PPTX: {filename}")
            text = extract_text_from_pptx(file_path)
        
        else:
            continue  
        
        if text.strip(): 
            clean_text = preprocess_text(text)
            chunk_sizes = [200, 500, 1000]
            overlaps = [0, 50, 100]
            
            for chunk_size in chunk_sizes:
                for overlap in overlaps:
                    chunks = chunk_text(clean_text, chunk_size, overlap)
                    all_chunks.extend(chunks) 
        
    return all_chunks

folder_path = "data/notes" 
chunks = process_notes_folder(folder_path)

# Save the chunks to a text file 
def save_chunks_to_file(chunks, output_folder="data/process_text"):
    """
    Saves the processed chunks to a text file in the specified folder.
    """
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    output_file = os.path.join(output_folder, "processed_chunks.txt")
    
    try:
        with open(output_file, "w", encoding="utf-8") as f:
            for i, chunk in enumerate(chunks):
                f.write(f"--- Chunk {i+1} ---\n{chunk}\n\n")
        print(f"Chunks successfully saved to {output_file}")
    except Exception as e:
        print(f"Error saving chunks to file: {e}")

save_chunks_to_file(chunks)
print("\nSample Chunks:")
for i, chunk in enumerate(chunks[:3]):  
    print(f"\n--- Chunk {i+1} ---\n{chunk}\n")
